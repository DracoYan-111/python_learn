from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "产科抗生素使用指南"
subtitle.text = "相关疾病、使用指针、诊断流程及送检病原"

# Slide 2: 相关疾病
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_2.shapes.title, slide_2.placeholders[1]
title.text = "相关疾病"
points = [
    "细菌感染",
    "支原体感染",
    "衣原体感染",
    "立克次体感染",
    "螺旋体感染",
    "真菌感染",
    "泌尿生殖系统感染",
    "术后感染",
    "绒毛膜羊膜炎"
]
for point in points:
    p = content.text_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(18)

# Slide 3: 使用指针
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_3.shapes.title, slide_3.placeholders[1]
title.text = "抗生素使用指针"
points = [
    "初步诊断为细菌性感染",
    "病原微生物检查确诊为细菌性感染",
    "选用正确的抗生素品种和给药方案",
    "进行病原微生物检查（细菌培养、药敏试验）"
]
for point in points:
    p = content.text_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(18)

# Slide 4: 诊断流程
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_4.shapes.title, slide_4.placeholders[1]
title.text = "一般诊断流程"
points = [
    "详细病史询问",
    "体格检查",
    "实验室检查（血常规、尿常规等）",
    "病原微生物培养"
]
for point in points:
    p = content.text_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(18)

# Slide 5: 送检病原
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_5.shapes.title, slide_5.placeholders[1]
title.text = "送检病原"
points = [
    "根据症状和疑似感染部位确定",
    "尿液",
    "血液",
    "生殖道分泌物"
]
for point in points:
    p = content.text_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(18)

# Save the presentation
pptx_path = '/mnt/data/Obstetric_Antibiotic_Use_Guide.pptx'
prs.save(pptx_path)

pptx_path
